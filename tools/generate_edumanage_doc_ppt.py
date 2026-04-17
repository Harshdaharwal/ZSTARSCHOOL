from __future__ import annotations

import datetime as _dt
from pathlib import Path


def _safe_rel(p: Path) -> str:
    try:
        return str(p.relative_to(Path.cwd()))
    except Exception:
        return str(p)


def _find_assets(repo_root: Path) -> dict:
    screenshots_dir = repo_root / "school-admin-react" / "src" / "screenshots"
    hero = repo_root / "school-admin-react" / "src" / "assets" / "hero.png"
    screenshots = []
    if screenshots_dir.exists():
        screenshots = sorted(
            [p for p in screenshots_dir.iterdir() if p.is_file() and p.suffix.lower() in {".png", ".jpg", ".jpeg"}],
            key=lambda p: p.name.lower(),
        )
    return {"hero": hero if hero.exists() else None, "screenshots": screenshots}


def build_docx(out_path: Path, assets: dict) -> None:
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Document defaults
    section = doc.sections[0]
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.8)
    section.left_margin = Inches(0.9)
    section.right_margin = Inches(0.9)

    title = doc.add_paragraph()
    run = title.add_run("EduManage")
    run.bold = True
    run.font.size = Pt(28)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    sub = doc.add_paragraph("School & Institute Management Software")
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER

    meta = doc.add_paragraph(
        f"Project Documentation\nDate: {_dt.date.today().strftime('%B %d, %Y')}"
    )
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_page_break()

    doc.add_heading("1. Introduction", level=1)
    doc.add_paragraph(
        "EduManage is a school/institute management software designed to streamline day-to-day academic and "
        "administrative workflows. It provides a unified interface for managing students, teachers, classes, "
        "attendance, fees, exams, marks, report cards, timetables, announcements, and staff salaries."
    )

    doc.add_heading("2. System Overview", level=1)
    doc.add_paragraph(
        "The application is built as a modern web dashboard. It supports role-based access for administrators and "
        "teachers, analytics on the dashboard, and a demo/offline-friendly mode using a local mock backend. "
        "Where configured, it can integrate with cloud services (Firebase) and/or a Google Apps Script (GAS) proxy."
    )

    doc.add_heading("3. User Roles & Access Control", level=1)
    doc.add_paragraph("EduManage supports the following roles:", style=None)
    doc.add_paragraph("Administrator (Admin)", style="List Bullet")
    doc.add_paragraph(
        "Full access to student records, fee management, exams/marks, settings, data snapshot exports, and administrative actions.",
        style="List Bullet",
    )
    doc.add_paragraph("Teacher", style="List Bullet")
    doc.add_paragraph(
        "Access to class-relevant dashboards, student results for assigned classes, posting class announcements, and attendance-related workflows.",
        style="List Bullet",
    )

    doc.add_paragraph(
        "Demo logins are provided for testing in local/mock mode. In production, authentication should be configured via Firebase or a backend."
    )

    doc.add_heading("4. Core Modules & Features", level=1)
    doc.add_paragraph(
        "EduManage is organized into clear modules so schools can manage each workflow independently while still "
        "keeping data consistent across the system (e.g., a student’s class/section is reused for attendance, fees, "
        "marks, report cards, and parent communication)."
    )
    modules = [
        ("Dashboard & Analytics", [
            "Overview cards (students, attendance, fees/collection where applicable).",
            "Charts for trends and class-wise distribution.",
            "Teacher: Today’s timetable panel based on schedule entries.",
        ]),
        ("Student Management", [
            "New student registration with photo upload (size-limited).",
            "Search, filter by class/section/academic year, pagination.",
            "Status management (Active/Inactive) and profile editing.",
        ]),
        ("Teacher Management", [
            "Add/update teacher profiles with photo and assignment (class/section).",
            "Delete teacher records (admin).",
        ]),
        ("Classes & Schedule", [
            "Create class-section records, assign class teacher, set room number.",
            "Create weekly schedule periods (day, period, subject, teacher, room, time slot).",
            "View schedule grouped by day and remove periods when needed.",
        ]),
        ("Attendance (Students & Teachers)", [
            "Mark or edit daily attendance by class/section/date.",
            "Student attendance codes: Present (P), Absent (A), Late (L).",
            "Attendance summaries per student/teacher.",
        ]),
        ("Fees Management", [
            "Add fee records (type, amount, due date, status).",
            "Quick views for All / Pending / Paid.",
            "Class fee structure setting (admin) for standardized fee rules by class.",
        ]),
        ("Exams, Marks & Results", [
            "Create exams (name, class, subject, date, max, pass).",
            "Enter marks and view marks list.",
            "Marks & Results matrix with filters and actions per student.",
            "Exports: CSV download, single PDF, class ZIP of PDFs (one per student).",
        ]),
        ("Report Cards", [
            "Generate single report card by student ID/name.",
            "Generate bulk report cards by class/section and print-friendly output.",
        ]),
        ("Timetables", [
            "Manage timetable entries for classes, tests, exams, and other events.",
            "Filter-friendly list with admin-only add/delete controls.",
        ]),
        ("Announcements", [
            "School-wide announcements with priority and optional attachment.",
            "Class announcements targeted to a specific class/section.",
        ]),
        ("Salary Management", [
            "Add salary records for teachers and staff.",
            "Filter by period (monthly/yearly/etc.) and status (paid/pending).",
            "Mark pending salaries as paid and maintain salary history.",
        ]),
        ("Settings, Audit & Backups", [
            "WhatsApp-style parent alerts (demo queue) with filters and logs.",
            "Audit log of recent actions for accountability.",
            "Admin data snapshot export as JSON for offline backup.",
        ]),
        ("Offline-Friendly Operation", [
            "When offline, write operations are queued in IndexedDB.",
            "When connectivity returns, queued operations sync automatically in order.",
        ]),
    ]
    for name, items in modules:
        doc.add_heading(name, level=2)
        for it in items:
            doc.add_paragraph(it, style="List Bullet")

    doc.add_heading("5. Typical Workflows (How Schools Use EduManage)", level=1)
    workflows = [
        (
            "Student admission / onboarding",
            [
                "Admin opens Students module and registers the student (name, parents, class/section, academic year, optional photo).",
                "Student becomes available immediately for attendance, fee records, and marks entry.",
                "For parent communication, admin/teacher can store Parent WhatsApp number (or fallback phone).",
            ],
        ),
        (
            "Daily attendance",
            [
                "Teacher/admin selects date + class + section and fetches student list.",
                "Mark Present/Absent/Late and save; edits can be done later via the Edit tab.",
                "Parents can be notified (demo WhatsApp queue) when a student is marked absent.",
            ],
        ),
        (
            "Fees collection",
            [
                "Admin records fee entries per student (monthly/annual/exam fee), sets due date and payment status.",
                "Use Pending/Paid views to track outstanding fees quickly.",
                "Optionally define a standard class fee structure so messages and reports can reference expected fees.",
            ],
        ),
        (
            "Exams, marks, and results",
            [
                "Admin creates an exam entry (class, subject, max marks, pass marks, date).",
                "Teacher/admin enters marks for students.",
                "Results module shows a class-wise matrix and supports per-student result view and PDF export.",
                "Admins can export CSV or a ZIP of PDFs for a class/section for offline sharing and printing.",
            ],
        ),
        (
            "Announcements and class notices",
            [
                "School-wide announcements: admin posts notices with priority and optional attachment.",
                "Class announcements: teacher/admin posts class-specific notice; parents of that class can be notified (demo WhatsApp queue).",
            ],
        ),
        (
            "Salary records (teachers and staff)",
            [
                "Admin adds monthly salary records with role/designation and amount.",
                "Filter by month/year and paid/pending; mark records as paid once processed.",
            ],
        ),
        (
            "Backup & audit",
            [
                "Admin downloads a JSON data snapshot for backup in demo mode.",
                "Audit log records key actions to help track changes and accountability.",
            ],
        ),
    ]
    for title, steps in workflows:
        doc.add_heading(title, level=2)
        for s in steps:
            doc.add_paragraph(s, style="List Number")

    doc.add_heading("6. Data Model (High-Level)", level=1)
    doc.add_paragraph(
        "The system centers around a few core entities. The UI forms map closely to these entities, which helps keep "
        "operations predictable and reduces duplicate entry."
    )
    entities = [
        ("Student", "ID, name, parents, class/section, academic year, contact info, status, optional photo."),
        ("Teacher", "ID, name, subject, contact details, class/section assignment, optional photo."),
        ("Class/Section", "Class number, section, class teacher, room number."),
        ("Schedule", "Weekly periods: day, period number, subject, teacher name/ID, room, time slot."),
        ("Attendance", "Daily attendance entries for students and teachers; summaries and percentage calculations."),
        ("Fee Record", "Fee type, amount, due date, paid date (if any), status, receipt/remarks."),
        ("Exam", "Exam metadata and per-student marks entries with grade and result."),
        ("Timetable Entry", "Class/test/exam events with date/day, time slot, room, notes."),
        ("Announcement", "School-wide or class-specific notices, priority, timestamps, optional attachment."),
        ("Salary Record", "Staff/teacher salary entries by month/year with paid/pending status."),
        ("Audit / Snapshot", "Recent actions log + data export for backup and reporting."),
    ]
    for name, desc in entities:
        doc.add_paragraph(f"{name}: {desc}", style="List Bullet")

    doc.add_heading("7. Architecture & Technology Stack", level=1)
    doc.add_paragraph("Frontend", style="List Bullet")
    doc.add_paragraph("React + Vite, React Router for navigation.", style="List Bullet 2")
    doc.add_paragraph("Chart.js for dashboards and analytics.", style="List Bullet 2")
    doc.add_paragraph("i18next for translations/localization support.", style="List Bullet 2")
    doc.add_paragraph("react-jss for styling.", style="List Bullet 2")
    doc.add_paragraph("Data & Backend Options", style="List Bullet")
    doc.add_paragraph(
        "Demo/mock backend for local usage and testing without cloud dependencies.", style="List Bullet 2"
    )
    doc.add_paragraph(
        "Firebase integration (Auth + Firestore/Realtime DB) when environment variables are configured.", style="List Bullet 2"
    )
    doc.add_paragraph(
        "Google Apps Script (GAS) proxy support when running in a GAS container environment.", style="List Bullet 2"
    )

    doc.add_heading("8. Setup & Running (Developer/Deployment Notes)", level=1)
    doc.add_paragraph("Prerequisites:", style="List Bullet")
    doc.add_paragraph("Node.js (for building/running the React app).", style="List Bullet 2")
    doc.add_paragraph("Typical commands:", style="List Bullet")
    doc.add_paragraph("Install dependencies: npm install", style="List Bullet 2")
    doc.add_paragraph("Run locally: npm run dev", style="List Bullet 2")
    doc.add_paragraph("Build: npm run build", style="List Bullet 2")
    doc.add_paragraph(
        "If Firebase is not configured, the app can run in demo/mock mode with demo credentials. "
        "For production, configure authentication and data storage, and replace demo notification behavior with real integrations."
    )

    doc.add_heading("9. Screens (From Project Screenshots)", level=1)
    shots = assets.get("screenshots") or []
    if shots:
        # Include several screenshots for a richer, 5–6 page document.
        for i, p in enumerate(shots[:6], start=1):
            doc.add_heading(f"Screen {i}", level=2)
            doc.add_paragraph(f"Source: {_safe_rel(p)}")
            try:
                doc.add_picture(str(p), width=Inches(6.2))
            except Exception:
                doc.add_paragraph("Unable to embed this image (unsupported or corrupted).")
    else:
        doc.add_paragraph("No screenshots were found in the repository.")

    doc.add_heading("10. Security, Privacy & Production Notes", level=1)
    for it in [
        "Avoid using demo credentials in production; use a secure identity provider and strong password policies.",
        "Do not store sensitive personal data beyond what is necessary for school operations; follow local compliance requirements.",
        "Treat exports (PDFs, ZIPs, CSV, JSON snapshots) as confidential; secure access and storage.",
        "Replace mock notification flows with audited, consent-based messaging via approved providers.",
    ]:
        doc.add_paragraph(it, style="List Bullet")

    doc.add_heading("11. Future Enhancements (Suggested)", level=1)
    for it in [
        "Integrate real WhatsApp/SMS/email provider for parent notifications.",
        "Add dedicated parent portal and student portal views.",
        "Stronger role management (multiple admin roles, permissions per module).",
        "Automated backups and scheduled reports.",
        "Improved reporting (custom templates, result analysis, dashboards per class/subject).",
    ]:
        doc.add_paragraph(it, style="List Bullet")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out_path))


def build_pptx(out_path: Path, assets: dict) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    hero = assets.get("hero")
    shots = assets.get("screenshots") or []
    shots = shots[:9]  # we need 9-ish images max (title + 8 content)

    def add_title(slide, title: str, subtitle: str | None = None):
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.2), Inches(1.2))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(36)
        run.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(18)
            p2.font.color.rgb = RGBColor(80, 80, 80)
            p2.level = 0

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    if hero and Path(hero).exists():
        slide.shapes.add_picture(str(hero), Inches(0), Inches(0), width=prs.slide_width)
        # Overlay title box for readability
        overlay = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE (avoid importing enum; stable id)
            Inches(0.6),
            Inches(5.6),
            Inches(12.2),
            Inches(1.6),
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
        overlay.fill.fore_color.brightness = 0.15
        overlay.line.color.rgb = RGBColor(255, 255, 255)
        tx = overlay.text_frame
        tx.clear()
        p = tx.paragraphs[0]
        r = p.add_run()
        r.text = "EduManage"
        r.font.size = Pt(40)
        r.font.bold = True
        p2 = tx.add_paragraph()
        p2.text = "School & Institute Management Software"
        p2.font.size = Pt(18)
    else:
        add_title(slide, "EduManage", "School & Institute Management Software")

    # Slide 2: Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Overview", "A unified platform for daily school operations")
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.2), Inches(5.2))
    tf = box.text_frame
    tf.clear()
    bullets = [
        "Student & teacher records",
        "Classes, schedules, and timetables",
        "Attendance (students & teachers)",
        "Fees, exams, marks & report cards",
        "Announcements + parent alerts (demo)",
        "Dashboard analytics + offline sync",
    ]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
    if shots:
        slide.shapes.add_picture(str(shots[0]), Inches(7.4), Inches(1.6), width=Inches(5.6))

    # Slide 3: Modules
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Key Modules", "Admin + teacher workflows")
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(12.0), Inches(5.6))
    tf = box.text_frame
    tf.clear()
    lines = [
        "Dashboard & analytics",
        "Students / Teachers / Classes",
        "Attendance and summaries",
        "Fees and class fee settings",
        "Exams, marks, results, exports (CSV/PDF/ZIP)",
        "Timetables and schedules",
        "Announcements (school-wide + class)",
        "Salary management",
        "Settings: audit log, data snapshot, WhatsApp queue (demo)",
    ]
    for i, t in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.level = 0
        p.font.size = Pt(18)

    # Slides 4-10: Screenshots (7 slides)
    screen_titles = [
        "Login & Role-Based Access",
        "Dashboard & Analytics",
        "Student Management",
        "Teacher Management",
        "Classes & Schedule",
        "Attendance / Results",
        "Settings / Announcements",
    ]
    # Use up to 7 more screenshots after the first one used on slide 2.
    screen_shots = shots[1:8]
    for idx in range(7):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, screen_titles[idx], "Screens from the EduManage project")

        if idx < len(screen_shots) and screen_shots[idx].exists():
            # Fit image into content area
            slide.shapes.add_picture(
                str(screen_shots[idx]),
                Inches(0.9),
                Inches(1.6),
                width=Inches(11.6),
            )
        else:
            box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.0))
            tf = box.text_frame
            tf.text = "Screenshot not available."
            tf.paragraphs[0].font.size = Pt(24)

        foot = slide.shapes.add_textbox(Inches(0.9), Inches(7.1), Inches(11.6), Inches(0.3))
        ft = foot.text_frame
        ft.clear()
        p = ft.paragraphs[0]
        p.text = "EduManage • School Management Software"
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.RIGHT
        p.font.color.rgb = RGBColor(90, 90, 90)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def build_client_docx(out_path: Path, assets: dict) -> None:
    """
    Client-ready product overview (brochure/proposal style).
    Keep it shareable: no demo credentials, no internal implementation details.
    """
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    section = doc.sections[0]
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.8)
    section.left_margin = Inches(0.9)
    section.right_margin = Inches(0.9)

    title = doc.add_paragraph()
    run = title.add_run("EduManage")
    run.bold = True
    run.font.size = Pt(30)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    sub = doc.add_paragraph("School & Institute Management Software")
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER

    tagline = doc.add_paragraph(
        "A modern, easy-to-use platform to manage students, staff, academics, fees, attendance, and communication — all in one place."
    )
    tagline.alignment = WD_ALIGN_PARAGRAPH.CENTER

    hero = assets.get("hero")
    if hero and Path(hero).exists():
        try:
            doc.add_picture(str(hero), width=Inches(6.2))
            doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
        except Exception:
            pass

    doc.add_heading("Product Summary", level=1)
    doc.add_paragraph(
        "EduManage helps schools and institutes reduce paperwork, improve accuracy, and save time by digitizing core workflows. "
        "It is designed for daily use by administrators and teachers, with clear modules for academics, operations, and reporting."
    )

    doc.add_heading("Who It’s For", level=1)
    for it in [
        "Schools (primary to senior secondary)",
        "Coaching centers and institutes",
        "School admins, principals, teachers, and office staff",
    ]:
        doc.add_paragraph(it, style="List Bullet")

    doc.add_heading("Key Benefits", level=1)
    for it in [
        "Centralized records: students, teachers, classes, fees, exams, and more",
        "Faster daily operations: attendance, announcements, and schedules",
        "Better visibility: dashboard analytics and reports",
        "Reduced errors: standardized data entry and consistent workflows",
        "Export-ready: CSV/PDF outputs where applicable for sharing and printing",
    ]:
        doc.add_paragraph(it, style="List Bullet")

    doc.add_heading("Core Features", level=1)
    features = [
        ("Student Management", "Register students, manage profiles, class/section, academic year, and status."),
        ("Teacher Management", "Add and manage teacher profiles, assignments, and contact details."),
        ("Classes & Schedule", "Maintain class/section, rooms, and weekly schedules (day/period/time/subject)."),
        ("Attendance", "Mark and edit daily attendance for students and teachers; view summaries."),
        ("Fees", "Record fee entries, track pending/paid status, and manage fee structure by class."),
        ("Exams & Marks", "Create exams, enter marks, and view results with filtering and exports."),
        ("Report Cards", "Generate single or bulk report cards for printing and sharing."),
        ("Timetables", "Add and manage timetable entries for class routines, tests, and exams."),
        ("Announcements", "Post school-wide and class-specific notices; attachments supported for announcements."),
        ("Salary Management", "Maintain teacher/staff salary records with paid/pending tracking."),
        ("Settings & Audit", "Admin tools including audit log and secure snapshot exports (where enabled)."),
    ]
    for name, desc in features:
        doc.add_paragraph(name, style="List Bullet")
        doc.add_paragraph(desc, style="List Bullet 2")

    doc.add_heading("Screens (Sample)", level=1)
    shots = assets.get("screenshots") or []
    if shots:
        for p in shots[:3]:
            try:
                doc.add_picture(str(p), width=Inches(6.2))
            except Exception:
                continue
    else:
        doc.add_paragraph("Screenshots available in the project build.")

    doc.add_heading("Implementation & Support (Typical)", level=1)
    for it in [
        "Requirement discussion and module selection",
        "Initial setup and configuration",
        "Data entry/import support (as required)",
        "Training for admin and teachers",
        "Ongoing support and updates",
    ]:
        doc.add_paragraph(it, style="List Bullet")

    doc.add_heading("Next Steps", level=1)
    doc.add_paragraph(
        "We can share a demo, finalize the modules you need, and provide an implementation plan based on your school size and workflow."
    )

    doc.add_paragraph()
    doc.add_paragraph("Contact / Proposal Details")
    doc.add_paragraph("School/Client Name: ___________________________")
    doc.add_paragraph("Prepared By: _________________________________")
    doc.add_paragraph("Date: ______________________________________")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out_path))


def main() -> int:
    repo_root = Path.cwd()
    assets = _find_assets(repo_root)

    out_dir = repo_root / "deliverables"
    docx_path = out_dir / "EduManage_Project_Documentation.docx"
    pptx_path = out_dir / "EduManage_Project_Presentation.pptx"
    client_docx_path = out_dir / "EduManage_Client_Product_Overview.docx"

    build_docx(docx_path, assets)
    build_pptx(pptx_path, assets)
    build_client_docx(client_docx_path, assets)

    print(f"Wrote: {_safe_rel(docx_path)}")
    print(f"Wrote: {_safe_rel(pptx_path)}")
    print(f"Wrote: {_safe_rel(client_docx_path)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
